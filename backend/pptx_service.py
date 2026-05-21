"""
pptx_service.py – converts text or ChatGPT JSON into a .pptx file.

Slide layout reference:
    0  – Title and Subtitle
    1  – Title and Content      ← used for every content slide
    2  – Section Header
    3  – Two Content
    4  – Comparison
    5  – Title Only
    6  – Blank
    7  – Content with Caption
    8  – Picture with Caption
"""

from pptx import Presentation
from io import BytesIO
import json


# ── Internal helpers ──────────────────────────────────────────────────────────

def _add_title_slide(prs: Presentation, title_text: str = "Slides") -> None:
    """Adds the opening title slide (layout 0)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = "Created by Sciensinsta"


def _add_content_slide(prs: Presentation, title: str, items: list[str]) -> None:
    """Adds a single content slide (layout 1) with bullet items."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    for item in items:
        item = item.strip()
        if item:                          # skip empty strings from split
            tf.add_paragraph().text = item


def _build_presentation(slides: list, is_gpt: bool) -> Presentation:
    """
    Builds and returns a Presentation object from a list of slides.

    For tool=1 (is_gpt=False):
        slides  – list of raw text paragraphs (split by '\\n\\n' upstream).
        Each paragraph is split on '.' to produce bullet points.
        Title is derived from the first sentence of the paragraph.

    For tool=2 (is_gpt=True):
        slides  – list of dicts with 'title' and 'content' keys (ChatGPT JSON).
        Content is split on '.' to produce bullet points.
    """
    prs = Presentation()
    _add_title_slide(prs)

    for slide in slides:
        if is_gpt:
            title = slide.get("title", "")
            content = slide.get("content", "")

            # content can be either a list of strings or a single string
            if isinstance(content, list):
                items = [s.strip() for s in content if str(s).strip()]
            else:
                items = [s.strip() for s in str(content).split(".") if s.strip()]
        else:
            sentences = [s.strip() for s in slide.split(".") if s.strip()]
            title = sentences[0] if sentences else ""
            items = sentences

        _add_content_slide(prs, title, items)

    return prs


# ── Public API ────────────────────────────────────────────────────────────────

def create_pptx(text_input: str, is_gpt: bool) -> BytesIO:
    """
    Generates a .pptx file entirely in memory and returns it as a BytesIO buffer.
    Nothing is written to disk.

    Parameters
    ----------
    text_input : str
        Plain text (tool=1) or JSON string (tool=2).
    is_gpt : bool
        False → split plain text by '\\n\\n' into slides.
        True  → parse text_input as a JSON list of {title, content} dicts.

    Returns
    -------
    BytesIO
        In-memory buffer positioned at byte 0, ready to stream.

    Raises
    ------
    ValueError
        If text_input is empty or (for tool=2) the JSON is invalid.
    """
    if not text_input or not text_input.strip():
        raise ValueError("text_input is empty")

    if is_gpt:
        try:
            slides = json.loads(text_input)
        except json.JSONDecodeError as exc:
            raise ValueError(f"Invalid JSON: {exc}") from exc
        if not isinstance(slides, list):
            raise ValueError("JSON must be a list of slide objects")
    else:
        slides = [s for s in text_input.split("\n\n") if s.strip()]
        if not slides:
            raise ValueError("No slide content found after splitting by double newline")

    prs = _build_presentation(slides, is_gpt)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)          # rewind so the caller can read from the start
    return buf
